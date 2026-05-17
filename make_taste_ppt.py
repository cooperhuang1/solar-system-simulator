'''
Scientific Taste PPT generator
Run: python make_taste_ppt.py
'''
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG     = RGBColor(0x08, 0x0C, 0x14)
PANEL  = RGBColor(0x10, 0x16, 0x24)
PANEL2 = RGBColor(0x16, 0x1E, 0x30)
BORDER = RGBColor(0x2A, 0x38, 0x52)
TEXT   = RGBColor(0xEE, 0xF2, 0xFA)
MUTED  = RGBColor(0x8A, 0x98, 0xB8)
FAINT  = RGBColor(0x4A, 0x58, 0x72)
GOLD   = RGBColor(0xF5, 0xC5, 0x42)
BLUE   = RGBColor(0x6A, 0xA8, 0xFF)
TEAL   = RGBColor(0x4E, 0xC9, 0xB0)
GREEN  = RGBColor(0x7E, 0xD3, 0x8B)
PINK   = RGBColor(0xE0, 0x7A, 0xA8)
ORANGE = RGBColor(0xF0, 0x9A, 0x4A)

W = Inches(13.33)
H = Inches(7.5)

LQ = '\u201c'  # left double curly quote
RQ = '\u201d'  # right double curly quote


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def box(slide, l, t, w, h, fill=None, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.line.width = lw
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=Pt(20), bold=False,
        color=TEXT, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Microsoft YaHei UI'
    return txb


def hline(slide, l, t, w, color=GOLD, h=Pt(3)):
    box(slide, l, t, w, h, fill=color, lw=Pt(0))


def vbar(slide, l, t, h, color=GOLD, w=Inches(0.16)):
    box(slide, l, t, w, h, fill=color, lw=Pt(0))


def page_header(slide, part_label, title, subtitle=None, accent=GOLD):
    vbar(slide, Inches(0), Inches(0), H, color=accent)
    box(slide, Inches(0.3), Inches(0.22), Inches(3.8), Inches(0.42),
        fill=PANEL2, line=BORDER, lw=Pt(0.5))
    txt(slide, part_label, Inches(0.42), Inches(0.24), Inches(3.6), Inches(0.38),
        size=Pt(12), color=accent)
    txt(slide, title, Inches(0.42), Inches(0.72), Inches(12.6), Inches(0.82),
        size=Pt(34), bold=True, color=TEXT)
    if subtitle:
        txt(slide, subtitle, Inches(0.42), Inches(1.5), Inches(12.6), Inches(0.45),
            size=Pt(16), color=MUTED, italic=True)
        hline(slide, Inches(0.42), Inches(1.95), Inches(12.5), color=BORDER, h=Pt(1))
    else:
        hline(slide, Inches(0.42), Inches(1.52), Inches(5), color=accent, h=Pt(3))


def quote_box(slide, l, t, w, h, quote, source, accent=GOLD):
    box(slide, l, t, w, h, fill=PANEL2, line=BORDER)
    hline(slide, l, t, w, color=accent, h=Pt(3))
    txt(slide, LQ + quote + RQ,
        l + Inches(0.3), t + Inches(0.22),
        w - Inches(0.6), h - Inches(0.7),
        size=Pt(16), color=TEXT, italic=True, wrap=True)
    txt(slide, '\u2014\u2014 ' + source,
        l + Inches(0.3), t + h - Inches(0.48),
        w - Inches(0.6), Inches(0.4),
        size=Pt(13), color=accent, align=PP_ALIGN.RIGHT)


def bullet_card(slide, lx, ty, w, h, title, items,
                accent=GOLD, title_size=Pt(17), body_size=Pt(14)):
    vbar(slide, lx, ty, h, color=accent)
    box(slide, lx + Inches(0.16), ty, w - Inches(0.16), h, fill=PANEL, line=BORDER)
    txt(slide, title, lx + Inches(0.34), ty + Inches(0.14),
        w - Inches(0.5), Inches(0.48), size=title_size, bold=True, color=accent)
    for i, item in enumerate(items):
        txt(slide, item, lx + Inches(0.38), ty + Inches(0.68 + i * 0.46),
            w - Inches(0.55), Inches(0.44), size=body_size, color=MUTED)


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    vbar(slide, Inches(0), Inches(0), H, color=GOLD)
    hline(slide, Inches(0.42), Inches(1.5), Inches(8.5), color=GOLD, h=Pt(3))
    txt(slide, '科学的隐形标尺', Inches(0.42), Inches(1.72), Inches(9.5), Inches(1.5),
        size=Pt(60), bold=True, color=TEXT)
    txt(slide, '论科学品味的重要性', Inches(0.42), Inches(3.1), Inches(9.5), Inches(0.9),
        size=Pt(36), color=GOLD)
    hline(slide, Inches(0.42), Inches(4.08), Inches(6), color=BORDER, h=Pt(1))
    txt(slide, '从' + LQ + '做研究' + RQ + '到' + LQ + '做卓越的研究' + RQ,
        Inches(0.42), Inches(4.2), Inches(9), Inches(0.6),
        size=Pt(20), color=MUTED, italic=True)
    txt(slide, '汇报人：黄俊鑫', Inches(0.42), Inches(5.1), Inches(5), Inches(0.5),
        size=Pt(18), color=MUTED)
    txt(slide, LQ, Inches(9.5), Inches(1.2), Inches(3.5), Inches(3.5),
        size=Pt(260), bold=True, color=RGBColor(0x1A, 0x24, 0x38))
    box(slide, Inches(0.42), Inches(6.7), Inches(12.5), Inches(0.55),
        fill=PANEL2, line=BORDER, lw=Pt(0.5))
    tags = ['科学哲学', '科研方法论', '学术成长', 'Scientific Taste']
    tx = 0.62
    for tag in tags:
        tw = len(tag) * 0.21 + 0.4
        box(slide, Inches(tx), Inches(6.78), Inches(tw), Inches(0.38),
            fill=PANEL, line=BORDER, lw=Pt(0.5))
        txt(slide, tag, Inches(tx + 0.1), Inches(6.8), Inches(tw), Inches(0.34),
            size=Pt(13), color=GOLD)
        tx += tw + 0.2


def slide_self_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '自我介绍', '关于我', accent=BLUE)
    box(slide, Inches(0.42), Inches(2.1), Inches(12.5), Inches(4.6),
        fill=PANEL, line=BORDER)
    txt(slide, '（此处留白，待填写个人信息）',
        Inches(0.42), Inches(4.0), Inches(12.5), Inches(0.8),
        size=Pt(22), color=FAINT, align=PP_ALIGN.CENTER)
    txt(slide, '姓名 / 院系 / 研究方向 / 代表性工作',
        Inches(0.42), Inches(4.7), Inches(12.5), Inches(0.5),
        size=Pt(16), color=FAINT, align=PP_ALIGN.CENTER)


def slide_soul_question(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第一部分  引言', '灵魂拷问', accent=GOLD)
    questions = [
        ('Q1', '为什么有的研究者极其勤奋，却只能产出平庸的成果？',   GOLD),
        ('Q2', '为什么有的研究者发文不多，但每一篇都能引领行业风向？', BLUE),
        ('Q3', '决定这两者差距的，究竟是什么？',                       GREEN),
    ]
    for i, (label, q, color) in enumerate(questions):
        ty = Inches(2.1 + i * 1.52)
        box(slide, Inches(0.42), ty, Inches(12.5), Inches(1.35), fill=PANEL, line=BORDER)
        box(slide, Inches(0.42), ty, Inches(0.72), Inches(1.35), fill=color, lw=Pt(0))
        txt(slide, label, Inches(0.5), ty + Inches(0.42),
            Inches(0.6), Inches(0.5), size=Pt(16), bold=True,
            color=RGBColor(0x08, 0x0C, 0x14), align=PP_ALIGN.CENTER)
        txt(slide, q, Inches(1.28), ty + Inches(0.3),
            Inches(11.4), Inches(0.75), size=Pt(22), color=TEXT)
    hline(slide, Inches(0.42), Inches(6.85), Inches(12.5), color=GOLD, h=Pt(2))
    txt(slide, '核心词：科学品味（Scientific Taste）',
        Inches(0.42), Inches(6.95), Inches(12.5), Inches(0.5),
        size=Pt(22), bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_masters_view(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第一部分  引言', '大师的视角', accent=GOLD)
    quote_box(slide, Inches(0.42), Inches(2.0), Inches(12.5), Inches(1.9),
              '科学研究中，提出一个问题往往比解决一个问题更重要。\n'
              '解决一个问题也许仅是一个数学上或实验上的技能而已，'
              '而提出新的问题、新的可能性，从新的角度去看旧的问题，'
              '却需要有创造性的想象力。',
              '阿尔伯特·爱因斯坦', accent=GOLD)
    quote_box(slide, Inches(0.42), Inches(4.1), Inches(12.5), Inches(1.55),
              '数学创造的本质不仅是逻辑，更是审美。'
              '数学家不是因为有用才研究数学，而是因为它美。',
              '亨利·庞加莱', accent=BLUE)
    box(slide, Inches(0.42), Inches(5.85), Inches(12.5), Inches(0.88),
        fill=PANEL2, line=BORDER)
    txt(slide, '结论：顶级科学家不仅是' + LQ + '解题高手' + RQ +
        '，更是具备极高品味的' + LQ + '问题挑选者' + RQ,
        Inches(0.7), Inches(5.95), Inches(12.0), Inches(0.65),
        size=Pt(19), bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_three_dimensions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第二部分  定义与内涵', '科学品味的三个维度', accent=TEAL)
    dims = [
        (GOLD,  'What  ·  选题的眼光',
         '能敏锐嗅到什么问题是根本性的、有价值的。\n'
         '不沉迷于修修补补的伪需求，而是直击领域核心矛盾。\n'
         '好的选题本身就已经成功了一半。'),
        (BLUE,  'How  ·  方法的审美',
         '追求简洁、优雅、直击本质的解决路径。\n'
         '摒弃繁琐和刻意复杂化——奥卡姆剃刀原理：\n'
         '如无必要，勿增实体。'),
        (GREEN, 'Where  ·  对未来的直觉',
         '在数据不充分时，对领域发展方向有超乎逻辑的判断力。\n'
         '这种直觉来自大量阅读、深度思考和长期积累，\n'
         '是品味最高级的体现。'),
    ]
    for i, (color, title, body) in enumerate(dims):
        lx = Inches(0.42 + i * 4.3)
        ty = Inches(2.0)
        vbar(slide, lx, ty, Inches(4.7), color=color, w=Inches(0.16))
        box(slide, lx + Inches(0.16), ty, Inches(4.0), Inches(4.7), fill=PANEL, line=BORDER)
        txt(slide, str(i + 1), lx + Inches(0.28), ty + Inches(0.12),
            Inches(0.6), Inches(0.7), size=Pt(42), bold=True, color=color)
        txt(slide, title, lx + Inches(0.28), ty + Inches(0.82),
            Inches(3.7), Inches(0.52), size=Pt(17), bold=True, color=color)
        hline(slide, lx + Inches(0.28), ty + Inches(1.36), Inches(3.6), color=BORDER, h=Pt(1))
        txt(slide, body, lx + Inches(0.28), ty + Inches(1.5),
            Inches(3.7), Inches(3.0), size=Pt(14), color=MUTED)


def slide_compare(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第二部分  定义与内涵', '高品味 vs. 低品味', accent=TEAL)
    box(slide, Inches(0.42), Inches(2.0), Inches(6.0), Inches(4.7), fill=PANEL, line=BORDER)
    hline(slide, Inches(0.42), Inches(2.0), Inches(6.0), color=PINK, h=Pt(4))
    txt(slide, '低品味  Low Taste', Inches(0.62), Inches(2.18),
        Inches(5.6), Inches(0.52), size=Pt(20), bold=True, color=PINK)
    for i, item in enumerate([
        '盲目追逐热点，跟风灌水',
        '用复杂的高射炮打蚊子',
        '为发 paper 而生造概念',
        '沉迷于边际改进，回避核心难题',
        '研究结论缺乏普遍意义',
    ]):
        txt(slide, '\u2717  ' + item, Inches(0.62), Inches(2.85 + i * 0.7),
            Inches(5.6), Inches(0.6), size=Pt(15), color=MUTED)
    box(slide, Inches(6.9), Inches(2.0), Inches(6.0), Inches(4.7), fill=PANEL, line=BORDER)
    hline(slide, Inches(6.9), Inches(2.0), Inches(6.0), color=GREEN, h=Pt(4))
    txt(slide, '高品味  High Taste', Inches(7.1), Inches(2.18),
        Inches(5.6), Inches(0.52), size=Pt(20), bold=True, color=GREEN)
    for i, item in enumerate([
        '啃硬骨头，做奠基性工作',
        '用最精简的模型解释复杂现象',
        '寻找普适性规律和底层逻辑',
        '在无人区开辟新方向',
        '每篇工作都能引发领域思考',
    ]):
        txt(slide, '\u2713  ' + item, Inches(7.1), Inches(2.85 + i * 0.7),
            Inches(5.6), Inches(0.6), size=Pt(15), color=MUTED)
    txt(slide, 'VS', Inches(6.1), Inches(3.9), Inches(0.8), Inches(0.8),
        size=Pt(28), bold=True, color=FAINT, align=PP_ALIGN.CENTER)


def slide_ceiling(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第三部分  为什么科学品味至关重要', '1. 决定研究的天花板', accent=ORANGE)
    box(slide, Inches(0.42), Inches(2.0), Inches(12.5), Inches(1.1), fill=PANEL2, line=BORDER)
    hline(slide, Inches(0.42), Inches(2.0), Inches(12.5), color=GOLD, h=Pt(4))
    txt(slide, '努力决定了你能走多快，品味决定了你能走多高。',
        Inches(0.7), Inches(2.12), Inches(12.0), Inches(0.75),
        size=Pt(26), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    left_items = [
        '方向大于努力',
        '在错误的方向上狂奔，产出的是正确的废话',
        '完美执行一个平庸的课题，依然是平庸',
        '选题决定了工作的理论上限',
    ]
    right_items = [
        '品味是选方向的能力',
        '高品味让你在起跑线上就领先',
        '一个好问题的价值远超十篇好文章',
        '顶级期刊拒绝的不是执行，是选题',
    ]
    for i, item in enumerate(left_items):
        ty = Inches(3.3 + i * 0.82)
        box(slide, Inches(0.42), ty, Inches(0.5), Inches(0.65), fill=ORANGE, lw=Pt(0))
        txt(slide, item, Inches(1.05), ty + Inches(0.08), Inches(5.5), Inches(0.55),
            size=Pt(15), color=MUTED)
    for i, item in enumerate(right_items):
        ty = Inches(3.3 + i * 0.82)
        box(slide, Inches(6.9), ty, Inches(0.5), Inches(0.65), fill=GOLD, lw=Pt(0))
        txt(slide, item, Inches(7.53), ty + Inches(0.08), Inches(5.5), Inches(0.55),
            size=Pt(15), color=MUTED)


def slide_involution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第三部分  为什么科学品味至关重要', '2. 突破内卷，避免低水平重复', accent=ORANGE)
    box(slide, Inches(0.42), Inches(2.0), Inches(12.5), Inches(1.0), fill=PANEL2, line=BORDER)
    hline(slide, Inches(0.42), Inches(2.0), Inches(12.5), color=BLUE, h=Pt(4))
    txt(slide, '在文献爆炸的时代，低壁垒的填坑式研究极易被替代。',
        Inches(0.7), Inches(2.1), Inches(12.0), Inches(0.7),
        size=Pt(22), bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    cards = [
        (PINK,   '红海困境',
         ['每年数百万篇论文涌现',
          '同质化研究泛滥，引用率极低',
          'AI 工具正在加速这一趋势',
          '低品味研究的生命周期越来越短']),
        (GREEN,  '蓝海突破',
         ['范式转移（Paradigm Shift）来自品味',
          '真正的创新往往在无人区',
          '好品味让你看到别人看不到的空白',
          '开辟新方向比填满旧方向价值高百倍']),
        (GOLD,   '品味的护城河',
         ['技术工具越来越廉价普及',
          'AI 可以执行，但无法选题',
          '提出好问题的能力无法被替代',
          '品味是未来科研的核心竞争力']),
    ]
    for i, (color, title, items) in enumerate(cards):
        lx = Inches(0.42 + i * 4.3)
        bullet_card(slide, lx, Inches(3.2), Inches(4.1), Inches(3.95),
                    title, items, accent=color, title_size=Pt(16), body_size=Pt(13))


def slide_compass(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第三部分  为什么科学品味至关重要', '3. 应对不确定性的指南针', accent=ORANGE)
    quote_box(slide, Inches(0.42), Inches(2.0), Inches(12.5), Inches(1.35),
              '科研的本质是探索未知，前面没有路标。在面临多条技术路线选择、'
              '或者遇到重大挫折时，品味和直觉是支撑你做决策和坚持下去的唯一信仰。',
              '核心观点', accent=ORANGE)
    points = [
        (GOLD,  '没有地图的旅程',
         '科研不像考试有标准答案。每一个重要决策——选方向、换路线、放弃还是坚持——'
         '都需要在信息不完整的情况下做出判断。'),
        (BLUE,  '品味是决策的底层逻辑',
         '当实验数据相互矛盾、文献观点分歧时，是品味告诉你哪条路更值得走。'
         '这不是玄学，而是大量积累后形成的高级直觉。'),
        (GREEN, '穿越低谷的定力',
         '几乎所有重大突破在早期都遭遇过质疑和冷落。'
         '是对这个问题足够重要的坚信，让科学家在寒冬中坚持下去。'),
    ]
    for i, (color, title, body) in enumerate(points):
        ty = Inches(3.5 + i * 1.22)
        vbar(slide, Inches(0.42), ty, Inches(1.08), color=color)
        box(slide, Inches(0.58), ty, Inches(12.34), Inches(1.08), fill=PANEL, line=BORDER)
        txt(slide, title, Inches(0.78), ty + Inches(0.1),
            Inches(3.5), Inches(0.42), size=Pt(16), bold=True, color=color)
        txt(slide, body, Inches(0.78), ty + Inches(0.54),
            Inches(12.0), Inches(0.5), size=Pt(14), color=MUTED)


def slide_magnet(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第三部分  为什么科学品味至关重要', '4. 吸引顶尖资源的磁石', accent=ORANGE)
    box(slide, Inches(0.42), Inches(2.0), Inches(12.5), Inches(0.95), fill=PANEL2, line=BORDER)
    hline(slide, Inches(0.42), Inches(2.0), Inches(12.5), color=GOLD, h=Pt(4))
    txt(slide, '高品味的研究天然具有吸引力——它是一块磁石。',
        Inches(0.7), Inches(2.1), Inches(12.0), Inches(0.65),
        size=Pt(24), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    magnets = [
        (BLUE,   '吸引最优秀的学生',
         '顶尖学生在选导师时，首先看的是研究品味和方向，而非短期发文数量。\n'
         '高品味的课题组形成正向循环：好学生 → 好工作 → 更好的学生。'),
        (GREEN,  '吸引顶级同行合作',
         '真正有价值的合作来自相互欣赏。高品味的工作让同行主动找上门，\n'
         '而不是你去求人合作。合作的质量决定了工作的上限。'),
        (PINK,   '获得评审专家认同',
         '顶刊审稿人见过太多平庸的工作。一篇有品味的文章，\n'
         '在摘要阶段就能让审稿人感受到这个问题值得被解决。'),
        (ORANGE, '获取优质科研资源',
         '基金委、奖项评审、邀请报告——这些资源的分配，\n'
         '本质上是同行对你科学品味的投票。'),
    ]
    for i, (color, title, body) in enumerate(magnets):
        col = i % 2
        row = i // 2
        lx = Inches(0.42 + col * 6.42)
        ty = Inches(3.1 + row * 2.0)
        vbar(slide, lx, ty, Inches(1.82), color=color)
        box(slide, lx + Inches(0.16), ty, Inches(6.1), Inches(1.82), fill=PANEL, line=BORDER)
        txt(slide, title, lx + Inches(0.34), ty + Inches(0.12),
            Inches(5.7), Inches(0.45), size=Pt(16), bold=True, color=color)
        txt(slide, body, lx + Inches(0.34), ty + Inches(0.6),
            Inches(5.7), Inches(1.1), size=Pt(13), color=MUTED)


def slide_dirac(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第四部分  经典案例', '案例一：寻找优雅的保罗·狄拉克', accent=TEAL)
    box(slide, Inches(0.42), Inches(2.0), Inches(7.5), Inches(4.8), fill=PANEL, line=BORDER)
    hline(slide, Inches(0.42), Inches(2.0), Inches(7.5), color=TEAL, h=Pt(4))
    txt(slide, 'Paul Dirac  ·  1902-1984  ·  诺贝尔物理学奖', Inches(0.62), Inches(2.15),
        Inches(7.1), Inches(0.45), size=Pt(14), color=TEAL)
    txt(slide, '故事', Inches(0.62), Inches(2.62),
        Inches(7.1), Inches(0.42), size=Pt(18), bold=True, color=TEXT)
    txt(slide,
        '狄拉克坚信：物理方程必须具备数学美。\n\n'
        '1928年，他在尝试将量子力学与狭义相对论统一时，'
        '拒绝了所有在他看来丑陋的数学形式。\n\n'
        '他基于对公式对称美的极致追求，推导出了著名的狄拉克方程。\n\n'
        '这个方程在数学上要求存在一种与电子质量相同、电荷相反的粒子——'
        '即反物质。彼时没有任何实验证据，但狄拉克相信方程的美不会撒谎。\n\n'
        '4年后，正电子被实验发现，完全证实了他的预言。',
        Inches(0.62), Inches(3.1), Inches(7.1), Inches(3.5), size=Pt(14), color=MUTED)
    box(slide, Inches(8.2), Inches(2.0), Inches(4.7), Inches(4.8), fill=PANEL2, line=BORDER)
    hline(slide, Inches(8.2), Inches(2.0), Inches(4.7), color=GOLD, h=Pt(4))
    txt(slide, '启示', Inches(8.4), Inches(2.15),
        Inches(4.3), Inches(0.45), size=Pt(18), bold=True, color=GOLD)
    for i, (title, body) in enumerate([
        ('审美即真理', '对数学美的追求直接通向了物理真理，品味不是装饰，是工具。'),
        ('拒绝妥协',   '狄拉克拒绝了所有凑合能用的方案，坚持寻找优雅的解。'),
        ('超越实验',   '在没有实验证据时，品味和直觉给了他预言的勇气。'),
        ('品味的预测力', '真正优雅的理论往往比实验跑得更快，这是品味的最高境界。'),
    ]):
        ty = Inches(2.75 + i * 1.0)
        txt(slide, '\u25b8 ' + title, Inches(8.4), ty,
            Inches(4.3), Inches(0.38), size=Pt(14), bold=True, color=GOLD)
        txt(slide, body, Inches(8.4), ty + Inches(0.38),
            Inches(4.3), Inches(0.52), size=Pt(13), color=MUTED)


def slide_hinton(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第四部分  经典案例', '案例二：不随波逐流的杰弗里·辛顿', accent=TEAL)
    box(slide, Inches(0.42), Inches(2.0), Inches(7.5), Inches(4.8), fill=PANEL, line=BORDER)
    hline(slide, Inches(0.42), Inches(2.0), Inches(7.5), color=TEAL, h=Pt(4))
    txt(slide, 'Geoffrey Hinton  ·  2024年诺贝尔物理学奖', Inches(0.62), Inches(2.15),
        Inches(7.1), Inches(0.45), size=Pt(14), color=TEAL)
    txt(slide, '故事', Inches(0.62), Inches(2.62),
        Inches(7.1), Inches(0.42), size=Pt(18), bold=True, color=TEXT)
    txt(slide,
        '1980年代至2000年代，神经网络研究经历了长达数十年的 AI 寒冬。\n\n'
        '学术界主流认为：神经网络计算代价高、理论不完善、不如支持向量机优雅。'
        '大量研究者转向了其他方向。\n\n'
        'Hinton 基于极高的科学品味和直觉，坚信神经网络是通向真正智能的正确道路。'
        '他在经费匮乏、同行质疑的环境中坚持了几十年。\n\n'
        '2012年，他的学生用深度神经网络在 ImageNet 竞赛中以压倒性优势获胜，'
        '直接开启了深度学习时代，彻底改变了 AI 领域。',
        Inches(0.62), Inches(3.1), Inches(7.1), Inches(3.5), size=Pt(14), color=MUTED)
    box(slide, Inches(8.2), Inches(2.0), Inches(4.7), Inches(4.8), fill=PANEL2, line=BORDER)
    hline(slide, Inches(8.2), Inches(2.0), Inches(4.7), color=BLUE, h=Pt(4))
    txt(slide, '启示', Inches(8.4), Inches(2.15),
        Inches(4.3), Inches(0.45), size=Pt(18), bold=True, color=BLUE)
    for i, (title, body) in enumerate([
        ('穿越周期的定力', '好品味赋予人不盲从主流的勇气，在寒冬中坚持正确的方向。'),
        ('长期主义',       '真正重要的问题往往需要几十年才能被验证，品味是长期投资。'),
        ('逆向思维',       '当所有人都在逃离一个方向时，高品味者看到的是机会而非风险。'),
        ('信念的来源',     'Hinton 的坚持不是固执，而是基于对问题本质的深刻理解。'),
    ]):
        ty = Inches(2.75 + i * 1.0)
        txt(slide, '\u25b8 ' + title, Inches(8.4), ty,
            Inches(4.3), Inches(0.38), size=Pt(14), bold=True, color=BLUE)
        txt(slide, body, Inches(8.4), ty + Inches(0.38),
            Inches(4.3), Inches(0.52), size=Pt(13), color=MUTED)


def slide_born(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第五部分  如何培养科学品味', '科学品味是天生的吗？', accent=GREEN)
    box(slide, Inches(0.42), Inches(2.0), Inches(12.5), Inches(1.0), fill=PANEL2, line=BORDER)
    hline(slide, Inches(0.42), Inches(2.0), Inches(12.5), color=GREEN, h=Pt(4))
    txt(slide, '品味如同艺术鉴赏力——部分靠直觉，但很大程度上可以后天刻意训练。',
        Inches(0.7), Inches(2.1), Inches(12.0), Inches(0.7),
        size=Pt(22), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    for i, (color, title, items) in enumerate([
        (PINK,  '误区：天才论',
         ['有些人天生就有品味', '我没有那种直觉', '品味是不可教的', '这些想法是逃避训练的借口']),
        (GREEN, '真相：可训练',
         ['品味 = 大量输入 + 深度思考', '每一位大师都经历了漫长的积累',
          '刻意练习可以系统提升品味', '关键是找到正确的训练方法']),
        (GOLD,  '训练的本质',
         ['大量接触高品味的工作', '建立自己的评价标准',
          '在实践中不断校准直觉', '与高品味的人深度交流']),
    ]):
        lx = Inches(0.42 + i * 4.3)
        bullet_card(slide, lx, Inches(3.2), Inches(4.1), Inches(3.95),
                    title, items, accent=color, title_size=Pt(16), body_size=Pt(14))


def slide_strategy1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第五部分  如何培养科学品味', '策略一：向上社交，与大师对话', accent=GREEN)
    for i, (color, title, body) in enumerate([
        (GOLD,  '读原始文献，而非综述',
         '多读经典 Original Papers，看大师如何提出问题、如何思考。\n'
         '综述告诉你是什么，原始文献告诉你为什么这样想。\n'
         '重点不是记住结论，而是学习思维方式。'),
        (BLUE,  '观察顶尖学者如何评价研究',
         '多参加高水平学术讲座，注意顶尖学者在 Q&A 环节提什么问题。\n'
         '他们的问题往往直击要害，这本身就是品味的示范。\n'
         '学会区分好问题和坏问题。'),
        (TEAL,  '建立自己的品味档案',
         '记录下让你觉得这个工作真漂亮的文章，分析美在哪里。\n'
         '同样记录下让你觉得这个工作很平庸的文章，分析差在哪里。\n'
         '长期积累，形成自己的评价体系。'),
        (GREEN, '主动寻求高质量反馈',
         '把自己的想法讲给最挑剔的人听，而不是最友善的人。\n'
         '真正有价值的反馈来自那些品味比你高的人。\n'
         '被否定不是失败，是校准品味的机会。'),
    ]):
        ty = Inches(2.0 + i * 1.3)
        vbar(slide, Inches(0.42), ty, Inches(1.16), color=color)
        box(slide, Inches(0.58), ty, Inches(12.34), Inches(1.16), fill=PANEL, line=BORDER)
        txt(slide, title, Inches(0.78), ty + Inches(0.1),
            Inches(4.5), Inches(0.42), size=Pt(16), bold=True, color=color)
        txt(slide, body, Inches(0.78), ty + Inches(0.54),
            Inches(12.0), Inches(0.58), size=Pt(13), color=MUTED)


def slide_strategy2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第五部分  如何培养科学品味', '策略二：训练批判性思维', accent=GREEN)
    box(slide, Inches(0.42), Inches(2.0), Inches(12.5), Inches(0.9), fill=PANEL2, line=BORDER)
    hline(slide, Inches(0.42), Inches(2.0), Inches(12.5), color=BLUE, h=Pt(4))
    txt(slide, '建立自己的学术品判标准，敢于对顶刊上的平庸文章说不。',
        Inches(0.7), Inches(2.08), Inches(12.0), Inches(0.65),
        size=Pt(20), bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(slide, '读每篇文献时，问自己三个问题：',
        Inches(0.42), Inches(3.05), Inches(12.5), Inches(0.5),
        size=Pt(18), bold=True, color=TEXT)
    for i, (color, title, body) in enumerate([
        (GOLD,  'Q1  这个问题真的重要吗？',
         '如果这个问题被解决了，领域会发生什么变化？\n'
         '是根本性的推进，还是边际性的改进？\n'
         '五年后，这个问题还会被人记得吗？'),
        (BLUE,  'Q2  方法是最直接的吗？',
         '作者是否用了最简洁的方式解决问题？\n'
         '有没有更优雅的路径被忽视了？\n'
         '复杂性是必要的，还是为了显得高级？'),
        (GREEN, 'Q3  结论有普遍意义吗？',
         '这个结论只适用于特定场景，还是揭示了普适规律？\n'
         '能否推广到其他领域？\n'
         '这项工作改变了我们对某件事的基本认知吗？'),
    ]):
        lx = Inches(0.42 + i * 4.3)
        ty = Inches(3.65)
        vbar(slide, lx, ty, Inches(3.5), color=color)
        box(slide, lx + Inches(0.16), ty, Inches(4.0), Inches(3.5), fill=PANEL, line=BORDER)
        txt(slide, title, lx + Inches(0.3), ty + Inches(0.15),
            Inches(3.7), Inches(0.48), size=Pt(15), bold=True, color=color)
        txt(slide, body, lx + Inches(0.3), ty + Inches(0.7),
            Inches(3.7), Inches(2.6), size=Pt(13), color=MUTED)


def slide_strategy3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第五部分  如何培养科学品味', '策略三：留出发呆的时间', accent=GREEN)
    quote_box(slide, Inches(0.42), Inches(2.0), Inches(12.5), Inches(1.3),
              '切忌让自己成为一直在做实验的机器。保持一定的前瞻性思考，'
              '跳出细节，从宏观（Big Picture）审视自己的研究领域。',
              '核心建议', accent=GREEN)
    for i, (color, title, body) in enumerate([
        (GOLD,   'Step Back 的价值',
         '当你深陷细节时，很容易失去对为什么做这件事的感知。\n'
         '定期强迫自己退出来，问：我在做的这件事，在整个领域的地图上处于什么位置？'),
        (BLUE,   '如何有效发呆',
         '散步、洗澡、睡前——这些时间是大脑整合信息的黄金时段。\n'
         '不要用手机填满所有空白时间。很多重要的想法诞生于无所事事的时刻。'),
        (TEAL,   '建立宏观视野的习惯',
         '每月读一篇领域综述，每季度问自己：领域的核心矛盾是什么？\n'
         '哪些问题五年后会变得重要？我现在的工作是否在正确的轨道上？'),
        (GREEN,  '跨领域的滋养',
         '品味的提升往往来自意想不到的地方。\n'
         '读历史、哲学、艺术——这些看似无用的输入，\n'
         '会在某个时刻以你意想不到的方式影响你的科学判断。'),
    ]):
        col = i % 2
        row = i // 2
        lx = Inches(0.42 + col * 6.42)
        ty = Inches(3.45 + row * 1.95)
        vbar(slide, lx, ty, Inches(1.78), color=color)
        box(slide, lx + Inches(0.16), ty, Inches(6.1), Inches(1.78), fill=PANEL, line=BORDER)
        txt(slide, title, lx + Inches(0.34), ty + Inches(0.12),
            Inches(5.7), Inches(0.42), size=Pt(15), bold=True, color=color)
        txt(slide, body, lx + Inches(0.34), ty + Inches(0.58),
            Inches(5.7), Inches(1.1), size=Pt(13), color=MUTED)


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, '第六部分  结语', '总结', accent=GOLD)
    box(slide, Inches(0.42), Inches(2.0), Inches(12.5), Inches(1.05), fill=PANEL2, line=BORDER)
    hline(slide, Inches(0.42), Inches(2.0), Inches(12.5), color=GOLD, h=Pt(4))
    txt(slide, '技术和工具正变得越来越廉价——AI、高端仪器、算力，人人可得。',
        Inches(0.7), Inches(2.1), Inches(12.0), Inches(0.72),
        size=Pt(20), color=TEXT, align=PP_ALIGN.CENTER)
    for i, (color, title, body) in enumerate([
        (GOLD,  '唯一的护城河',
         '在未来，提出好问题的能力——即科学品味——将成为科研工作者唯一且最核心的护城河。\n'
         'AI 可以执行实验、分析数据、撰写论文，但它无法替代人类对什么值得研究的判断。'),
        (BLUE,  '从工匠到艺术家',
         '工匠精神是必要的，但不够。卓越的科研需要艺术家的眼光——\n'
         '对美的感知、对本质的追求、对未来的直觉。这正是科学品味的全部内涵。'),
        (GREEN, '品味是可以修炼的',
         '读经典、问好问题、与大师对话、留出思考的空间。\n'
         '每一次刻意的训练，都在悄悄提升你的隐形标尺。'),
        (TEAL,  '从今天开始',
         '下一次读文献时，不只问这篇文章说了什么，\n'
         '而是问这个问题值得被解决吗？这个方法足够优雅吗？'),
    ]):
        ty = Inches(3.2 + i * 1.05)
        vbar(slide, Inches(0.42), ty, Inches(0.92), color=color)
        box(slide, Inches(0.58), ty, Inches(12.34), Inches(0.92), fill=PANEL, line=BORDER)
        txt(slide, title, Inches(0.78), ty + Inches(0.06),
            Inches(3.0), Inches(0.38), size=Pt(15), bold=True, color=color)
        txt(slide, body, Inches(0.78), ty + Inches(0.46),
            Inches(12.0), Inches(0.42), size=Pt(13), color=MUTED)


def slide_qa(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    vbar(slide, Inches(0), Inches(0), H, color=GOLD)
    hline(slide, Inches(0.42), Inches(2.2), Inches(8), color=GOLD, h=Pt(3))
    txt(slide, 'Q & A', Inches(0.42), Inches(2.45), Inches(9), Inches(1.4),
        size=Pt(72), bold=True, color=TEXT)
    txt(slide, '欢迎提问交流', Inches(0.42), Inches(3.75), Inches(9), Inches(0.7),
        size=Pt(28), color=GOLD)
    quote_box(slide, Inches(0.42), Inches(4.6), Inches(9.5), Inches(1.55),
              '愿我们都能从一个科研的工匠，蜕变为具备高级品味的科学艺术家。',
              '共勉', accent=GOLD)
    txt(slide, '感谢聆听', Inches(0.42), Inches(6.35), Inches(9), Inches(0.6),
        size=Pt(20), color=MUTED)
    txt(slide, '?', Inches(10.0), Inches(1.5), Inches(3.0), Inches(4.5),
        size=Pt(280), bold=True, color=RGBColor(0x14, 0x1E, 0x32),
        align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    slide_cover(prs)
    slide_self_intro(prs)
    slide_soul_question(prs)
    slide_masters_view(prs)
    slide_three_dimensions(prs)
    slide_compare(prs)
    slide_ceiling(prs)
    slide_involution(prs)
    slide_compass(prs)
    slide_magnet(prs)
    slide_dirac(prs)
    slide_hinton(prs)
    slide_born(prs)
    slide_strategy1(prs)
    slide_strategy2(prs)
    slide_strategy3(prs)
    slide_conclusion(prs)
    slide_qa(prs)
    out = os.path.normpath(
        os.path.join(os.path.dirname(os.path.abspath(__file__)),
                     '科学品味_Scientific_Taste.pptx'))
    prs.save(out)
    print('saved: ' + out)


if __name__ == '__main__':
    build()
