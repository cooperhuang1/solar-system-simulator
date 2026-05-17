"""
生成太阳系模拟器项目介绍 PPT
运行: python -m solar_system_app.make_ppt
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── 颜色主题 ──────────────────────────────────────────────────────
BG      = RGBColor(0x05, 0x07, 0x0B)
PANEL   = RGBColor(0x0B, 0x0F, 0x16)
PANEL2  = RGBColor(0x0F, 0x14, 0x1D)
BORDER  = RGBColor(0x26, 0x31, 0x42)
TEXT    = RGBColor(0xE6, 0xED, 0xF7)
MUTED   = RGBColor(0x87, 0x93, 0xA6)
FAINT   = RGBColor(0x4E, 0x5B, 0x6F)
ACCENT  = RGBColor(0x77, 0xA7, 0xFF)
WARNING = RGBColor(0xF2, 0xC1, 0x66)
SUN     = RGBColor(0xFF, 0xD1, 0x66)
TEAL    = RGBColor(0x79, 0xD2, 0xD8)
GREEN   = RGBColor(0xA8, 0xD8, 0x8A)
PINK    = RGBColor(0xD8, 0x8A, 0xA8)

W = Inches(13.33)
H = Inches(7.5)

PLANET_COLORS = [
    RGBColor(0xA8,0xA3,0x9B), RGBColor(0xD9,0xB3,0x6C),
    RGBColor(0x4A,0x90,0xFF), RGBColor(0xD8,0x6A,0x4D),
    RGBColor(0xD9,0xA6,0x6E), RGBColor(0xE4,0xC2,0x77),
    RGBColor(0x79,0xD2,0xD8), RGBColor(0x55,0x78,0xFF),
]


# ── 基础绘制工具 ──────────────────────────────────────────────────

def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def box(slide, left, top, width, height, fill_color=None, line_color=None, lw=Pt(0.75)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.width = lw
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=Pt(20), bold=False, color=TEXT, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Microsoft YaHei UI"
    return txb


def hline(slide, left, top, width, color=ACCENT, h=Pt(2.5)):
    box(slide, left, top, width, h, fill_color=color, lw=Pt(0))


def left_bar(slide, left, top, height, color=ACCENT, w=Inches(0.14)):
    box(slide, left, top, w, height, fill_color=color, lw=Pt(0))


def page_header(slide, title, subtitle=None):
    """统一页眉：左竖条 + 大标题 + 可选副标题 + 分隔线"""
    box(slide, Inches(0), Inches(0), Inches(0.18), H, fill_color=ACCENT, lw=Pt(0))
    txt(slide, title, Inches(0.45), Inches(0.28), Inches(12.5), Inches(0.75),
        size=Pt(36), bold=True, color=TEXT)
    if subtitle:
        txt(slide, subtitle, Inches(0.45), Inches(0.95), Inches(12.5), Inches(0.45),
            size=Pt(16), color=MUTED)
        hline(slide, Inches(0.45), Inches(1.38), Inches(12.5), color=BORDER, h=Pt(1))
    else:
        hline(slide, Inches(0.45), Inches(1.05), Inches(5), color=ACCENT, h=Pt(3))


# 第1页：封面
def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # 左侧蓝色竖条
    box(slide, Inches(0), Inches(0), Inches(0.22), H, fill_color=ACCENT, lw=Pt(0))
    # 顶部装饰线
    hline(slide, Inches(0.55), Inches(1.65), Inches(7.5), color=ACCENT, h=Pt(3))

    # 主标题
    txt(slide, "太阳系模拟器", Inches(0.55), Inches(1.85), Inches(9), Inches(1.6),
        size=Pt(64), bold=True, color=TEXT)
    # 英文
    txt(slide, "Solar System Simulator", Inches(0.55), Inches(3.35), Inches(9), Inches(0.8),
        size=Pt(30), color=ACCENT)
    # 一句话描述
    txt(slide, "基于 Python + Tkinter 的实时三维太阳系轨道仿真程序",
        Inches(0.55), Inches(4.2), Inches(9.5), Inches(0.6),
        size=Pt(18), color=MUTED)
    # 技术标签行
    tags = ["Python 3", "Tkinter Canvas", "开普勒力学", "3D 透视投影", "实时仿真"]
    tx = 0.55
    for tag in tags:
        w = len(tag) * 0.22 + 0.35
        box(slide, Inches(tx), Inches(5.0), Inches(w), Inches(0.42),
            fill_color=PANEL2, line_color=BORDER, lw=Pt(0.75))
        txt(slide, tag, Inches(tx + 0.1), Inches(5.02), Inches(w - 0.1), Inches(0.38),
            size=Pt(13), color=ACCENT)
        tx += w + 0.18

    # 右侧太阳装饰
    halo = slide.shapes.add_shape(9, Inches(9.8), Inches(1.7), Inches(3.5), Inches(3.5))
    halo.fill.background()
    halo.line.color.rgb = RGBColor(0x5A, 0x42, 0x1F)
    halo.line.width = Pt(1.5)
    sun = slide.shapes.add_shape(9, Inches(10.15), Inches(2.05), Inches(2.8), Inches(2.8))
    sun.fill.solid(); sun.fill.fore_color.rgb = SUN
    sun.line.fill.background()
    # 内光晕
    inner = slide.shapes.add_shape(9, Inches(10.5), Inches(2.4), Inches(2.1), Inches(2.1))
    inner.fill.solid(); inner.fill.fore_color.rgb = RGBColor(0xFF, 0xE8, 0xA0)
    inner.line.fill.background()

    # 底部信息栏
    box(slide, Inches(0.22), Inches(6.85), W - Inches(0.22), Inches(0.65),
        fill_color=PANEL2, lw=Pt(0))
    txt(slide, "Python 课程汇报  ·  2026  ·  太阳系轨道仿真项目",
        Inches(0.55), Inches(6.92), Inches(8), Inches(0.45),
        size=Pt(14), color=MUTED)
    txt(slide, "6 模块  ·  ~800 行代码",
        Inches(10.5), Inches(6.92), Inches(2.5), Inches(0.45),
        size=Pt(14), color=FAINT, align=PP_ALIGN.RIGHT)


# 第2页：项目概览
def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, "项目概览",
                "一个用纯 Python 实现的可交互太阳系实时仿真程序，无需任何第三方图形库")

    cards = [
        (ACCENT,   "🎯  项目目标",
         "直观展示八大行星的轨道运动规律，将开普勒轨道力学与实时可视化结合，"
         "让用户能够以任意视角观察太阳系，并通过时间加速感受行星运动的尺度差异。"),
        (WARNING,  "🛠  技术选型",
         "Python 3 标准库 Tkinter Canvas 作为唯一渲染后端，无需安装 PyGame / OpenGL 等依赖。"
         "开普勒方程牛顿迭代求解行星位置，透视投影矩阵实现 3D 效果。"),
        (TEAL,     "📦  项目规模",
         "6 个模块，约 800 行代码。分为数据层、物理引擎、相机系统、渲染引擎、应用主控五大层次，"
         "架构清晰，各模块职责单一，便于维护和扩展。"),
        (GREEN,    "✨  核心亮点",
         "· 真实 J2000.0 轨道根数，启动自动同步当前日期\n"
         "· 3D 透视旋转 + 焦点跟随任意天体\n"
         "· 时间倍率 1~900 天/秒，支持暂停/加速\n"
         "· 完整侧边栏遥测面板 + 底部控制按钮"),
    ]

    for i, (color, title, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        lx = Inches(0.45 + col * 6.42)
        ty = Inches(1.55 + row * 2.82)
        left_bar(slide, lx, ty, Inches(2.6), color=color)
        box(slide, lx + Inches(0.14), ty, Inches(6.1), Inches(2.6),
            fill_color=PANEL, line_color=BORDER)
        txt(slide, title, lx + Inches(0.32), ty + Inches(0.18),
            Inches(5.7), Inches(0.5), size=Pt(17), bold=True, color=color)
        txt(slide, body, lx + Inches(0.32), ty + Inches(0.72),
            Inches(5.7), Inches(1.75), size=Pt(14), color=MUTED)

# --- 第3页：模块架构 ---
def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, "模块架构",
                "6 个模块分层协作，职责单一，数据单向流动：catalog → models → orbit → camera → renderer → application")

    modules = [
        ("catalog.py",    "数据层",     WARNING,
         "存储八大行星完整数据：轨道六根数、物理参数（半径/质量/重力/自转/温度）、"
         "中英文名称及科普描述。所有数据以 frozen dataclass 形式提供，运行时不可变。"),
        ("models.py",     "数据模型",   ACCENT,
         "定义核心数据结构：CelestialBody（天体）、OrbitalElements（轨道根数）、"
         "RenderBody（渲染缓存）、SimulationState（仿真状态）。使用 Python dataclass，类型安全。"),
        ("orbit.py",      "物理引擎",   TEAL,
         "开普勒方程牛顿迭代求解偏近点角，将轨道根数转换为日心黄道三维坐标。"
         "提供轨道折线（260点）和运动尾迹（48点）采样，内置 LRU 缓存加速重复计算。"),
        ("camera.py",     "相机系统",   GREEN,
         "维护 yaw/pitch/zoom 三个相机参数，提供偏航俯仰旋转和透视投影。"
         "投影公式：perspective = 920/(920+z)，实现近大远小的真实空间感。"),
        ("renderer.py",   "渲染引擎",   PINK,
         "6 层 Canvas 绘制流水线：背景星空 → 参考网格 → 轨道椭圆 → 运动尾迹 → 天体本体 → UI面板。"
         "侧边栏宽 392px，底部控制栏高 82px，所有 UI 元素支持点击命中检测。"),
        ("application.py","应用主控",   WARNING,
         "16ms 定时器驱动游戏循环，绑定全部键鼠事件。管理仿真状态机：暂停/运行、"
         "时间倍率（1~900天/秒）、焦点跟随、标签/轨迹开关、天体选中逻辑。"),
    ]

    for i, (name, layer, color, desc) in enumerate(modules):
        col = i % 2
        row = i // 2
        lx = Inches(0.45 + col * 6.42)
        ty = Inches(1.55 + row * 1.92)
        left_bar(slide, lx, ty, Inches(1.78), color=color)
        box(slide, lx + Inches(0.14), ty, Inches(6.1), Inches(1.78),
            fill_color=PANEL, line_color=BORDER)
        txt(slide, name, lx + Inches(0.32), ty + Inches(0.12),
            Inches(3.2), Inches(0.45), size=Pt(16), bold=True, color=color)
        txt(slide, layer, lx + Inches(4.5), ty + Inches(0.14),
            Inches(1.6), Inches(0.38), size=Pt(13), color=MUTED, align=PP_ALIGN.RIGHT)
        txt(slide, desc, lx + Inches(0.32), ty + Inches(0.58),
            Inches(5.7), Inches(1.1), size=Pt(13), color=MUTED)


# 第4页：物理引擎
def slide_physics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, "物理引擎  ·  开普勒轨道力学",
                "以 J2000.0（2000年1月1日12:00 UTC）为历元，用六个轨道根数完整描述行星运动")

    # ── 左侧：求解流程 ──
    steps = [
        (ACCENT,  "① 六个轨道根数",
         "半长轴 a（AU）、离心率 e、轨道倾角 i、升交点经度 Ω、近日点辐角 ω、历元平近点角 M₀ + 周期 T"),
        (WARNING, "② 计算平近点角 M",
         "M = M₀ + 2π × t / T，其中 t 为距 J2000 的天数，随仿真时间线性递增"),
        (TEAL,    "③ 求解开普勒方程",
         "超越方程 E − e·sin(E) = M，用牛顿迭代法求解偏近点角 E，迭代 8 次精度达 1e-10 rad"),
        (GREEN,   "④ 轨道平面直角坐标",
         "x = a(cos E − e)，  y = a√(1−e²)·sin E，得到轨道平面内的位置向量"),
        (PINK,    "⑤ 旋转至日心黄道系",
         "依次绕 z 轴旋转近日点辐角 ω、绕 x 轴旋转轨道倾角 i、绕 z 轴旋转升交点经度 Ω"),
        (WARNING, "⑥ 屏幕距离缩放",
         "紧凑模式：r_screen = r_AU^0.66 × 132（压缩外行星距离便于展示）\n"
         "线性模式：r_screen = r_AU × 58（保持真实比例）"),
    ]

    for i, (color, title, body) in enumerate(steps):
        ty = Inches(1.55 + i * 0.97)
        box(slide, Inches(0.45), ty, Inches(0.14), Inches(0.84),
            fill_color=color, lw=Pt(0))
        txt(slide, title, Inches(0.72), ty + Inches(0.04),
            Inches(5.8), Inches(0.38), size=Pt(14), bold=True, color=color)
        txt(slide, body, Inches(0.72), ty + Inches(0.42),
            Inches(5.8), Inches(0.5), size=Pt(12), color=MUTED)

    # ── 右侧：轨道根数数据表 ──
    box(slide, Inches(7.1), Inches(1.55), Inches(5.95), Inches(5.75),
        fill_color=PANEL, line_color=BORDER)
    txt(slide, "八大行星轨道根数", Inches(7.3), Inches(1.68),
        Inches(5.5), Inches(0.45), size=Pt(16), bold=True, color=TEXT)
    hline(slide, Inches(7.3), Inches(2.15), Inches(5.5), color=BORDER, h=Pt(1))

    headers = ["行星", "a (AU)", "e", "i (°)", "T (天)"]
    col_x   = [7.32, 8.52, 9.52, 10.42, 11.38]
    col_w   = [1.1,  0.9,  0.8,  0.85,  1.1]
    ty = Inches(2.22)
    for j, h in enumerate(headers):
        txt(slide, h, Inches(col_x[j]), ty, Inches(col_w[j]), Inches(0.38),
            size=Pt(13), bold=True, color=ACCENT)

    planets = [
        ("水星", "0.3871", "0.2056", "7.005",   "87.97"),
        ("金星", "0.7233", "0.0068", "3.395",   "224.7"),
        ("地球", "1.0000", "0.0167", "0.000",   "365.3"),
        ("火星", "1.5237", "0.0934", "1.850",   "687.0"),
        ("木星", "5.2028", "0.0489", "1.303",  "4332.6"),
        ("土星", "9.5388", "0.0565", "2.485", "10759.2"),
        ("天王星","19.191", "0.0472", "0.773", "30685.4"),
        ("海王星","30.061", "0.0086", "1.770", "60189.0"),
    ]
    for k, (row_data, pcolor) in enumerate(zip(planets, PLANET_COLORS)):
        ry = Inches(2.65 + k * 0.575)
        bg = PANEL2 if k % 2 == 0 else PANEL
        box(slide, Inches(7.1), ry - Inches(0.04), Inches(5.95), Inches(0.575),
            fill_color=bg, lw=Pt(0))
        # 行星色点
        dot = slide.shapes.add_shape(9, Inches(7.32), ry + Inches(0.1),
                                     Inches(0.22), Inches(0.22))
        dot.fill.solid(); dot.fill.fore_color.rgb = pcolor
        dot.line.fill.background()
        for j, val in enumerate(row_data):
            c = TEXT if j == 0 else MUTED
            ox = Inches(col_x[j]) + (Inches(0.28) if j == 0 else 0)
            txt(slide, val, ox, ry, Inches(col_w[j]), Inches(0.45),
                size=Pt(13), color=c)


# 第5页：渲染系统
def slide_rendering(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, "渲染系统  ·  相机与投影",
                "纯 Tkinter Canvas 实现 3D 透视渲染，每帧约 16ms，6 层绘制流水线")

    # ── 左列：相机系统 ──
    box(slide, Inches(0.45), Inches(1.55), Inches(6.1), Inches(5.65),
        fill_color=PANEL, line_color=BORDER)
    hline(slide, Inches(0.45), Inches(1.55), Inches(6.1), color=ACCENT, h=Pt(3))
    txt(slide, "相机系统  camera.py", Inches(0.65), Inches(1.68),
        Inches(5.7), Inches(0.48), size=Pt(17), bold=True, color=ACCENT)
    hline(slide, Inches(0.65), Inches(2.18), Inches(5.7), color=BORDER, h=Pt(1))

    cam_items = [
        (WARNING, "偏航 / 俯仰旋转",
         "鼠标拖拽驱动 yaw（左右）和 pitch（上下）。\n"
         "俯仰角限制在 12°~86° 之间，防止视角翻转穿越极点。\n"
         "每像素旋转量：yaw×0.006 rad，pitch×0.004 rad。"),
        (TEAL,    "透视投影",
         "投影距离 d=920 单位，perspective = d/(d+z)。\n"
         "屏幕坐标 = center + world_xy × zoom × perspective。\n"
         "z 值同时用于天体深度排序，保证正确遮挡关系。"),
        (GREEN,   "缩放与焦点跟随",
         "滚轮调节 zoom ∈ [0.38, 3.2]，每格缩放 ×1.1 或 ×0.91。\n"
         "follow_focus 模式：投影前将世界坐标减去目标天体位置，\n"
         "实现镜头锁定跟随任意行星。"),
    ]
    for i, (color, title, body) in enumerate(cam_items):
        ty = Inches(2.3 + i * 1.6)
        left_bar(slide, Inches(0.65), ty, Inches(1.38), color=color, w=Inches(0.1))
        txt(slide, title, Inches(0.88), ty + Inches(0.05),
            Inches(5.1), Inches(0.42), size=Pt(15), bold=True, color=color)
        txt(slide, body, Inches(0.88), ty + Inches(0.5),
            Inches(5.1), Inches(0.95), size=Pt(13), color=MUTED)

    # ── 右列：渲染层次 ──
    box(slide, Inches(6.9), Inches(1.55), Inches(6.1), Inches(5.65),
        fill_color=PANEL, line_color=BORDER)
    hline(slide, Inches(6.9), Inches(1.55), Inches(6.1), color=WARNING, h=Pt(3))
    txt(slide, "渲染流水线  renderer.py", Inches(7.1), Inches(1.68),
        Inches(5.7), Inches(0.48), size=Pt(17), bold=True, color=WARNING)
    hline(slide, Inches(7.1), Inches(2.18), Inches(5.7), color=BORDER, h=Pt(1))

    layers = [
        (RGBColor(0x26,0x31,0x42), "① 背景 + 星空",
         "深色背景填充，520 颗随机星点（固定种子 20260515），\n三种亮度等级模拟真实星空深度感。"),
        (RGBColor(0x17,0x22,0x33), "② 参考网格",
         "96px 间距深色格线覆盖视口，外加内边框矩形，\n提供空间参考坐标感。"),
        (RGBColor(0x4A,0x90,0xFF), "③ 轨道椭圆 + 尾迹",
         "每条轨道 260 个采样点绘制折线，选中行星高亮蓝色。\n尾迹 48 点，覆盖约 22% 轨道周期，颜色渐隐。"),
        (RGBColor(0xD9,0xA6,0x6E), "④ 天体本体",
         "按 z 深度排序后绘制，保证正确遮挡。太阳绘制光晕圆，\n土星额外绘制倾斜光环椭圆，选中天体显示瞄准框。"),
        (RGBColor(0xE6,0xED,0xF7), "⑤ UI 面板",
         "右侧 392px 侧边栏（遥测数据 + 天体目录），\n底部 82px 控制栏（8个功能按钮），悬停名称提示。"),
    ]
    for i, (color, title, body) in enumerate(layers):
        ty = Inches(2.3 + i * 0.97)
        box(slide, Inches(6.9), ty, Inches(0.18), Inches(0.84),
            fill_color=color, lw=Pt(0))
        txt(slide, title, Inches(7.18), ty + Inches(0.04),
            Inches(5.6), Inches(0.38), size=Pt(14), bold=True, color=TEXT)
        txt(slide, body, Inches(7.18), ty + Inches(0.44),
            Inches(5.6), Inches(0.48), size=Pt(12), color=MUTED)


# 第6页：交互功能
def slide_interaction(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, "交互功能",
                "完整的键鼠交互体系，所有操作均有即时视觉反馈")

    features = [
        (ACCENT,  "🖱  鼠标操作",
         ["左键拖拽：旋转视角（偏航 + 俯仰）",
          "滚轮：缩放视图（0.38× ~ 3.2×，每格 ×1.1）",
          "左键单击天体：选中并切换侧边栏信息",
          "左键单击底部按钮：触发对应功能",
          "鼠标悬停天体：显示名称浮动提示框"]),
        (WARNING, "⌨  键盘快捷键",
         ["Space：暂停 / 继续仿真",
          "+  /  −：时间倍率 ×1.45 / ÷1.45（范围 1~900 天/秒）",
          "← →：按目录顺序切换选中天体",
          "0：聚焦太阳，1~8：直接跳转对应行星",
          "L 标签  T 轨迹  F 跟随  R 重置视角"]),
        (TEAL,    "📊  侧边栏遥测面板",
         ["TELEMETRY 标题区：天体名称、中文名、类型",
          "物理参数：半径、质量、重力、自转周期、温度",
          "轨道参数：半长轴、离心率、轨道倾角、周期",
          "OBJECT CATALOG：可点击的天体列表（含色点）",
          "CONTROLS：快捷键速查表"]),
        (GREEN,   "⏱  时间与状态系统",
         ["启动时自动读取系统时间，同步当前 UTC 日期",
          "仿真时间以 J2000.0 为历元，单位为天",
          "时间倍率默认 18 天/秒，可调至 900 天/秒",
          "状态栏实时显示：运行状态 / 日期 / 倍率 / 比例模式",
          "紧凑/线性两种轨道比例模式随时切换"]),
    ]

    for i, (color, title, items) in enumerate(features):
        col = i % 2
        row = i // 2
        lx = Inches(0.45 + col * 6.42)
        ty = Inches(1.55 + row * 2.82)
        left_bar(slide, lx, ty, Inches(2.68), color=color)
        box(slide, lx + Inches(0.14), ty, Inches(6.1), Inches(2.68),
            fill_color=PANEL, line_color=BORDER)
        txt(slide, title, lx + Inches(0.32), ty + Inches(0.15),
            Inches(5.7), Inches(0.48), size=Pt(17), bold=True, color=color)
        for j, item in enumerate(items):
            txt(slide, f"·  {item}", lx + Inches(0.35), ty + Inches(0.7 + j * 0.4),
                Inches(5.6), Inches(0.38), size=Pt(13), color=MUTED)


# 第7页：界面说明 & 运行方式
def slide_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, "界面说明  ·  运行方式",
                "窗口默认 1440×860，最小 1120×720，所有布局自适应窗口大小")

    # 启动命令框
    box(slide, Inches(0.45), Inches(1.55), Inches(12.6), Inches(0.92),
        fill_color=PANEL2, line_color=BORDER)
    txt(slide, "▶   python -m solar_system_app",
        Inches(0.75), Inches(1.65), Inches(11.5), Inches(0.65),
        size=Pt(24), bold=True, color=SUN)

    # 四个界面区域说明
    regions = [
        (ACCENT,  "主视口（左上大区域）",
         "3D 透视星空场景，实时渲染八大行星位置、轨道椭圆和运动尾迹。\n"
         "支持鼠标拖拽旋转、滚轮缩放、点击选中天体。\n"
         "选中天体显示蓝色高亮轨道和瞄准框，跟随模式下镜头锁定目标。"),
        (WARNING, "侧边栏（右侧 392px）",
         "上半部分：TELEMETRY 遥测面板，显示选中天体的完整物理参数和轨道根数。\n"
         "下半部分：OBJECT CATALOG 天体目录，点击任意行可切换选中目标。\n"
         "底部：CONTROLS 快捷键速查表。"),
        (TEAL,    "底部控制栏（高 82px）",
         "8 个功能按钮：PAUSE/RUN、RATE−、RATE+、FOLLOW ON/OFF、\n"
         "LABELS ON/OFF、TRAILS ON/OFF、COMPACT/LINEAR、RESET VIEW。\n"
         "按钮状态与仿真状态实时同步，点击即时生效。"),
        (GREEN,   "状态栏（左上角面板）",
         "实时显示四项信息：RUNNING/PAUSED 运行状态、\n"
         "DATE 仿真日期（UTC 格式）、RATE 时间倍率（天/秒）、\n"
         "当前比例模式（COMPACT SCALE / LINEAR AU SCALE）。"),
    ]

    for i, (color, title, body) in enumerate(regions):
        col = i % 2
        row = i // 2
        lx = Inches(0.45 + col * 6.42)
        ty = Inches(2.65 + row * 2.3)
        left_bar(slide, lx, ty, Inches(2.16), color=color)
        box(slide, lx + Inches(0.14), ty, Inches(6.1), Inches(2.16),
            fill_color=PANEL, line_color=BORDER)
        txt(slide, title, lx + Inches(0.32), ty + Inches(0.14),
            Inches(5.7), Inches(0.45), size=Pt(16), bold=True, color=color)
        txt(slide, body, lx + Inches(0.32), ty + Inches(0.62),
            Inches(5.7), Inches(1.42), size=Pt(13), color=MUTED)


# 第8页：总结
def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    page_header(slide, "总结与展望",
                "项目完整实现了从轨道力学计算到实时 3D 可视化的全链路，具备良好的扩展基础")

    points = [
        (ACCENT,  "物理准确性",
         "采用 NASA JPL 发布的 J2000.0 轨道根数，开普勒方程 8 次牛顿迭代，"
         "位置精度优于 1e-10 rad。仿真时间与真实日历完全对应，可复现任意历史日期的行星构型。"),
        (WARNING, "工程架构",
         "严格分层的 MVC 架构：catalog/models（数据）→ orbit（物理）→ camera（视图）"
         "→ renderer（渲染）→ application（控制）。模块间依赖单向，便于单独测试和替换。"),
        (TEAL,    "交互体验",
         "16ms 帧循环（约 60fps）保证流畅动画。完整键鼠交互，时间倍率跨越 900 倍，"
         "支持焦点跟随、轨迹显示、两种比例模式，适合课堂演示和自主探索。"),
        (GREEN,   "可扩展方向",
         "① 添加月球、矮行星（冥王星）、小行星带\n"
         "② 引入 VSOP87 高精度行星理论提升精度\n"
         "③ 迁移至 Pygame / OpenGL 支持更大场景和粒子效果\n"
         "④ 增加行星合相、冲日等天文事件自动检测"),
    ]

    for i, (color, title, body) in enumerate(points):
        ty = Inches(1.55 + i * 1.38)
        left_bar(slide, Inches(0.45), ty, Inches(1.24), color=color)
        box(slide, Inches(0.59), ty, Inches(12.46), Inches(1.24),
            fill_color=PANEL, line_color=BORDER)
        txt(slide, title, Inches(0.82), ty + Inches(0.1),
            Inches(2.8), Inches(0.45), size=Pt(17), bold=True, color=color)
        txt(slide, body, Inches(0.82), ty + Inches(0.58),
            Inches(12.0), Inches(0.62), size=Pt(14), color=MUTED)

    # 底部感谢
    hline(slide, Inches(0.45), Inches(7.0), Inches(12.6), color=BORDER, h=Pt(1))
    txt(slide, "感谢观看  /  Thank You",
        Inches(0), Inches(7.1), W, Inches(0.5),
        size=Pt(20), bold=True, color=MUTED, align=PP_ALIGN.CENTER)


# 主入口
def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_cover(prs)
    slide_overview(prs)
    slide_architecture(prs)
    slide_physics(prs)
    slide_rendering(prs)
    slide_interaction(prs)
    slide_demo(prs)
    slide_summary(prs)

    out = os.path.normpath(
        os.path.join(os.path.dirname(__file__), "..", "太阳系模拟器_项目介绍_v2.pptx")
    )
    prs.save(out)
    print(f"✅  已保存：{out}")


if __name__ == "__main__":
    build()